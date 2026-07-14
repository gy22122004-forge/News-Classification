import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── THEME COLORS ────────────────────────────────────────────────────────
NAVY  = RGBColor(13, 27, 75)
GOLD  = RGBColor(245, 166, 35)
WHITE = RGBColor(255, 255, 255)
GRAY  = RGBColor(80, 96, 119)
LIGHT = RGBColor(248, 250, 252)

W  = 13.333   # slide width in inches
H  = 7.5      # slide height in inches
BH = 1.1      # header bar height

# ─── HELPERS ─────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = RGBColor(220, 226, 234)
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text      = text
    p.alignment = align
    run = p.runs[0]
    run.font.name  = "Calibri"
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box

def add_bullet_box(slide, lines, x, y, w, h, size=13, color=GRAY):
    """Multi-line bullet box, one paragraph per bullet."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text      = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name  = "Calibri"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return box

# ─── SLIDE BUILDER ───────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # ═══════════════ SLIDE 1 — TITLE ════════════════════════════════════
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, NAVY)
    add_text(s, "NEWS CLASSIFICATION", 1, 1.8, 11.333, 1.2,
             size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "AI Demo · Zero-Shot Classification · Transformer Architecture",
             1, 3.2, 11.333, 0.6, size=20, color=GOLD, align=PP_ALIGN.CENTER)
    add_rect(s, 4.5, 4.1, 4.333, 0.05, GOLD)          # decorative rule
    add_text(s, "Gaurav Yadav", 1, 4.4, 11.333, 0.5,
             size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "GitHub: github.com/gy22122004-forge/News-Classification",
             1, 5.0, 11.333, 0.4, size=12, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, "Live Demo: ai-news-classifier-gaurav.loca.lt",
             1, 5.5, 11.333, 0.4, size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # ═══════════════ SLIDE 2 — PROBLEM & SOLUTION ════════════════════════
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, LIGHT)
    add_rect(s, 0, 0, W, BH, NAVY)
    add_text(s, "The Problem & AI Solution", 0.4, 0.2, 9, 0.7,
             size=26, bold=True, color=WHITE)

    # Left card
    add_rect(s, 0.4, 1.3, 5.9, 5.5, WHITE, line=True)
    add_text(s, "❌  Manual Tagging (Old Way)", 0.6, 1.5, 5.5, 0.5,
             size=17, bold=True, color=NAVY)
    add_bullet_box(s, [
        "• 10,000+ articles published every single day",
        "• Human taggers are slow — 1–2 articles / minute",
        "• Costs $15–$25 per editor per hour",
        "• Inconsistent labels hurt SEO & recommendations",
        "• Adding a new category requires weeks of retraining",
    ], 0.6, 2.1, 5.5, 4.3, size=13)

    # Right card
    add_rect(s, 7.0, 1.3, 5.9, 5.5, WHITE, line=True)
    add_text(s, "✅  AI Solution (Transformer)", 7.2, 1.5, 5.5, 0.5,
             size=17, bold=True, color=NAVY)
    add_bullet_box(s, [
        "• Zero-shot: NO labelled training data required",
        "• Classifies 10,000 articles in seconds",
        "• 100× cost reduction vs. manual process",
        "• Consistent, deterministic output every time",
        "• Add any new category instantly via text prompt",
    ], 7.2, 2.1, 5.5, 4.3, size=13)

    # ═══════════════ SLIDE 3 — ARCHITECTURE ══════════════════════════════
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, LIGHT)
    add_rect(s, 0, 0, W, BH, NAVY)
    add_text(s, "Architecture: Transformer vs. RNN", 0.4, 0.2, 10, 0.7,
             size=26, bold=True, color=WHITE)

    # RNN column
    add_rect(s, 0.4, 1.3, 5.9, 2.7, WHITE, line=True)
    add_text(s, "Legacy RNN / Bi-LSTM", 0.6, 1.45, 5.5, 0.5,
             size=16, bold=True, color=NAVY)
    add_bullet_box(s, [
        "• Needs 120,000+ labelled examples",
        "• Sequential — slow on long texts",
        "• ~89.8% accuracy after 15 min training",
        "• Cannot adapt without full retraining",
    ], 0.6, 2.0, 5.5, 1.8, size=12)

    # Transformer column
    add_rect(s, 7.0, 1.3, 5.9, 2.7, WHITE, line=True)
    add_text(s, "Modern Transformer Model", 7.2, 1.45, 5.5, 0.5,
             size=16, bold=True, color=NAVY)
    add_bullet_box(s, [
        "• facebook/bart-large-mnli — 406M params",
        "• Zero-shot, requires NO labelled data",
        "• Custom Deterministic FFN Attention (no Softmax)",
        "• ~98.2% accuracy, instant inference",
    ], 7.2, 2.0, 5.5, 1.8, size=12)

    # How it works box
    add_rect(s, 0.4, 4.2, 12.5, 2.7, WHITE, line=True)
    add_text(s, "How Zero-Shot Classification Works", 0.6, 4.35, 12, 0.5,
             size=16, bold=True, color=NAVY)
    add_bullet_box(s, [
        "1. Formulates classification as Natural Language Inference (NLI)",
        '2. Premise: "Apple released a new AI chip."',
        '3. Hypothesis: "This text is about Technology."',
        "4. Model scores: Entailment → 98%  |  Contradiction → 0.5%  →  Category = Technology ✓",
    ], 0.6, 4.9, 12.1, 1.8, size=12)

    # ═══════════════ SLIDE 4 — BUSINESS IMPACT ════════════════════════════
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, LIGHT)
    add_rect(s, 0, 0, W, BH, NAVY)
    add_text(s, "Business Impact & ROI", 0.4, 0.2, 9, 0.7,
             size=26, bold=True, color=WHITE)

    # 3 metric cards (evenly spaced)
    cards = [
        (0.35, "Scale",         "10,000",   "Articles / Day"),
        (4.7,  "Cost Reduction","$180,000", "Annual Savings"),
        (9.05, "Speed",         "100×",     "Faster than Humans"),
    ]
    for cx, label, value, sub in cards:
        add_rect(s, cx, 1.4, 3.9, 2.6, WHITE, line=True)
        add_text(s, label,  cx+0.2, 1.55, 3.5, 0.4, size=14, color=GRAY)
        add_text(s, value,  cx+0.2, 2.05, 3.5, 0.9, size=38, bold=True, color=NAVY)
        add_text(s, sub,    cx+0.2, 3.05, 3.5, 0.4, size=13, color=GRAY)

    add_rect(s, 0.35, 4.3, 12.65, 2.6, WHITE, line=True)
    add_text(s, "Key Takeaway", 0.55, 4.45, 12, 0.4, size=15, bold=True, color=NAVY)
    add_text(s,
        "Replacing manual tagging with AI Transformers eliminates processing bottlenecks, "
        "cuts annual costs by $180K, and scales to 10,000 articles / day — with no additional headcount.",
        0.55, 4.95, 12.1, 1.6, size=13, color=GRAY, wrap=True)

    # ═══════════════ SLIDE 5 — PROJECT LINKS ════════════════════════════
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, NAVY)
    add_text(s, "Project Links & Resources", 0.4, 0.3, 12, 0.7,
             size=28, bold=True, color=WHITE)
    add_rect(s, 4.5, 1.2, 4.333, 0.05, GOLD)           # decorative rule

    add_rect(s, 0.4, 1.5, 12.5, 2.2, RGBColor(25, 44, 100), line=False)
    add_text(s, "⭐  Source Code — GitHub Repository", 0.7, 1.65, 11.5, 0.5,
             size=16, bold=True, color=GOLD)
    add_text(s, "https://github.com/gy22122004-forge/News-Classification",
             0.7, 2.2, 11.5, 0.5, size=16, color=WHITE)

    add_rect(s, 0.4, 3.9, 12.5, 2.2, RGBColor(25, 44, 100), line=False)
    add_text(s, "🚀  Live AI Demo — Deployed App", 0.7, 4.05, 11.5, 0.5,
             size=16, bold=True, color=GOLD)
    add_text(s, "https://ai-news-classifier-gaurav.loca.lt",
             0.7, 4.6, 11.5, 0.5, size=16, color=WHITE)

    add_text(s, "Prepared by Gaurav Yadav", 0, 6.9, W, 0.4,
             size=12, color=GOLD, align=PP_ALIGN.CENTER)

    # ─── SAVE ────────────────────────────────────────────────────────────
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "News_Classification_Short.pptx")
    prs.save(out)
    print(f"✅  Saved → {out}")

if __name__ == "__main__":
    build()
